"""Rich Office export helpers for markdown reports.

The app still keeps the lightweight DOCX/PPTX generators from the Galaxy VOC source,
but this module adds report-focused, in-memory Word/PPTX exports for benchmarking.
"""
from __future__ import annotations

import io
import re
from datetime import datetime
from typing import Any


def strip_md(text: str) -> str:
    text = str(text or "")
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1", text)
    return text.strip()


def split_sections(markdown: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    title = "요약"
    lines: list[str] = []
    for raw in (markdown or "").splitlines():
        line = raw.rstrip()
        if line.startswith("## "):
            if lines:
                sections.append((title, lines))
            title = strip_md(line.lstrip("# ").strip())
            lines = []
        elif line.startswith("# ") and not sections and not lines:
            title = strip_md(line.lstrip("# ").strip())
        else:
            lines.append(line)
    if lines:
        sections.append((title, lines))
    return sections or [("보고서", (markdown or "").splitlines())]


def markdown_to_docx_bytes(markdown: str, metadata: dict[str, Any] | None = None) -> bytes:
    try:
        from docx import Document
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Pt, RGBColor
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-docx가 필요합니다. requirements.txt 설치 후 다시 실행하세요.") from exc

    meta = metadata or {}
    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Malgun Gothic"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    styles["Normal"].font.size = Pt(10.5)

    title = meta.get("title") or "AI VOC 벤치마킹 보고서"
    p = doc.add_paragraph()
    r = p.add_run(title)
    r.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(30, 64, 175)
    if meta.get("subtitle"):
        p = doc.add_paragraph(meta["subtitle"])
        p.style = doc.styles["Subtitle"] if "Subtitle" in [s.name for s in doc.styles] else doc.styles["Normal"]
    doc.add_paragraph(meta.get("meta_line") or f"생성일: {datetime.now():%Y-%m-%d %H:%M}")
    doc.add_paragraph("")

    def add_heading(text: str, level: int = 1):
        para = doc.add_heading(strip_md(text), level=min(max(level, 1), 3))
        for run in para.runs:
            run.font.name = "Malgun Gothic"
            run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")

    for raw in (markdown or "").splitlines():
        line = raw.rstrip()
        if not line:
            continue
        if line.startswith("# "):
            add_heading(line[2:], 1)
        elif line.startswith("## "):
            add_heading(line[3:], 2)
        elif line.startswith("### "):
            add_heading(line[4:], 3)
        elif line.lstrip().startswith(("- ", "* ")):
            doc.add_paragraph(strip_md(line.lstrip()[2:]), style="List Bullet")
        elif re.match(r"^\d+\.\s+", line.strip()):
            doc.add_paragraph(strip_md(line.strip()), style="List Number")
        elif line.startswith("|") and line.endswith("|"):
            # Simple markdown tables are preserved as monospaced text; robust table parsing is fragile.
            p = doc.add_paragraph(strip_md(line))
            for run in p.runs:
                run.font.name = "Consolas"
                run.font.size = Pt(9)
        elif line.startswith("```"):
            continue
        else:
            doc.add_paragraph(strip_md(line))

    # footer-like core properties
    core = doc.core_properties
    core.author = meta.get("author") or "AI VOC Benchmark Studio Pro"
    core.title = title
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def markdown_to_pptx_bytes(markdown: str, metadata: dict[str, Any] | None = None, max_slides: int = 12) -> bytes:
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-pptx가 필요합니다. requirements.txt 설치 후 다시 실행하세요.") from exc

    meta = metadata or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(30, 64, 175)
    NAVY = RGBColor(17, 24, 39)
    GRAY = RGBColor(75, 85, 99)
    LIGHT = RGBColor(248, 250, 252)
    WHITE = RGBColor(255, 255, 255)

    def add_title(slide, text, y=0.35, size=26, color=BLUE):
        box = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = strip_md(text)[:90]
        run.font.name = "맑은 고딕"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
        return box

    def add_body(slide, lines, y=1.15):
        box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.85), Inches(5.6))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        kept = []
        for raw in lines:
            line = strip_md(raw)
            if not line or line.startswith("|") or line.startswith("---") or line.startswith("```"):
                continue
            line = re.sub(r"^[-*]\s+", "", line)
            if len(line) > 150:
                line = line[:147] + "..."
            kept.append(line)
            if len(kept) >= 8:
                break
        if not kept:
            kept = ["상세 내용은 Word/Markdown 보고서를 참고하세요."]
        for i, line in enumerate(kept):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ""
            p.level = 0
            run = p.add_run()
            run.text = f"• {line}"
            run.font.name = "맑은 고딕"
            run.font.size = Pt(15 if len(kept) <= 5 else 12)
            run.font.color.rgb = GRAY
            p.space_after = Pt(8)

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(1, Inches(0.75), Inches(0.75), Inches(11.85), Inches(1.35))
    band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.color.rgb = BLUE
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.05), Inches(11.3), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = meta.get("title") or "AI VOC 벤치마킹 보고서"; r.font.name = "맑은 고딕"; r.font.bold = True; r.font.size = Pt(30); r.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(1.25), Inches(2.65), Inches(10.8), Inches(0.7))
    p = sub.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = meta.get("subtitle") or f"생성일: {datetime.now():%Y-%m-%d %H:%M}"; r.font.name = "맑은 고딕"; r.font.size = Pt(18); r.font.color.rgb = NAVY

    sections = split_sections(markdown)
    for title, lines in sections[: max_slides - 1]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = WHITE
        add_title(slide, title)
        sep = slide.shapes.add_shape(1, Inches(0.58), Inches(0.95), Inches(12.1), Inches(0.04))
        sep.fill.solid(); sep.fill.fore_color.rgb = BLUE; sep.line.fill.background()
        add_body(slide, lines)
        foot = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.25))
        p = foot.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = "AI VOC Benchmark Studio Pro"; r.font.name = "맑은 고딕"; r.font.size = Pt(9); r.font.color.rgb = GRAY

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
